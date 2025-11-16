"""
PowerPoint CEO Deck Generator
Creates presentation-ready executive report
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Non-interactive backend
from io import BytesIO
import logging

logger = logging.getLogger(__name__)


class PowerPointDeckGenerator:
    """
    Generates PowerPoint presentation from executive report data.
    Simple, visual, CEO-ready.
    """
    
    def generate_deck(self, report_data: dict) -> BytesIO:
        """
        Generate PowerPoint deck from report data.
        Returns BytesIO buffer with .pptx file.
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Title
        self._add_title_slide(prs, report_data)
        
        # Slide 2: Executive Summary
        self._add_executive_summary(prs, report_data)
        
        # Slide 3: Top Bottlenecks
        self._add_bottlenecks_slide(prs, report_data)
        
        # Slide 4: People Analysis
        self._add_people_slide(prs, report_data)
        
        # Slide 5: Recommendations
        self._add_recommendations_slide(prs, report_data)
        
        # Save to buffer
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        return buffer
    
    def _add_title_slide(self, prs, data):
        """Title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = data['executive_summary']['headline']
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 51, 102)  # Navy
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Bottleneck Analysis Report - {data['period']}"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Date
        date_box = slide.shapes.add_textbox(Inches(3), Inches(6), Inches(4), Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = f"Generated: {data['generated_at'][:10]}"
        date_para = date_frame.paragraphs[0]
        date_para.font.size = Pt(14)
        date_para.font.color.rgb = RGBColor(150, 150, 150)
        date_para.alignment = PP_ALIGN.CENTER
    
    def _add_executive_summary(self, prs, data):
        """Executive summary slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        # TLDR
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = data['executive_summary']['tldr']
        p.font.size = Pt(20)
        p.line_spacing = 1.5
        
        # Key numbers
        numbers = data['executive_summary']['key_numbers']
        
        # Add metrics boxes
        y_pos = 5
        metrics = [
            (f"${numbers['total_blocked_value']/1000000:.1f}M", "Total Blocked Value"),
            (f"{numbers['people_bottlenecks']}", "People Bottlenecks"),
            (f"${numbers['recovery_potential']/1000000:.1f}M", "Recovery Potential")
        ]
        
        for i, (value, label) in enumerate(metrics):
            x_pos = 1 + (i * 2.5)
            
            # Value
            value_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos), Inches(2), Inches(0.6))
            value_tf = value_box.text_frame
            value_p = value_tf.paragraphs[0]
            value_p.text = value
            value_p.font.size = Pt(32)
            value_p.font.bold = True
            value_p.font.color.rgb = RGBColor(220, 38, 38)  # Red
            value_p.alignment = PP_ALIGN.CENTER
            
            # Label
            label_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos + 0.7), Inches(2), Inches(0.4))
            label_tf = label_box.text_frame
            label_p = label_tf.paragraphs[0]
            label_p.text = label
            label_p.font.size = Pt(12)
            label_p.font.color.rgb = RGBColor(100, 100, 100)
            label_p.alignment = PP_ALIGN.CENTER
    
    def _add_bottlenecks_slide(self, prs, data):
        """Bottlenecks slide with chart"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Top Bottlenecks"
        
        # Get findings
        findings = data.get('key_findings', [])
        if not findings:
            return
        
        # Add text explanation
        text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for i, finding in enumerate(findings[:3], 1):
            p = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
            p.text = f"{i}. {finding['title']}: {finding['simple_explanation']}"
            p.font.size = Pt(16)
            p.space_after = Pt(12)
            p.level = 0
    
    def _add_people_slide(self, prs, data):
        """People bottlenecks slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Who's Overloaded"
        
        # Explanation
        finding = next((f for f in data.get('key_findings', []) if f['title'] == "Who's Overloaded"), None)
        if finding:
            text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
            tf = text_box.text_frame
            tf.text = finding['simple_explanation']
            tf.paragraphs[0].font.size = Pt(18)
            
            # Show top 3 people
            if finding.get('data'):
                y_pos = 4
                for person in finding['data'][:3]:
                    name_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(3), Inches(0.5))
                    name_box.text_frame.text = person.get('person', 'Unknown')
                    name_box.text_frame.paragraphs[0].font.size = Pt(16)
                    name_box.text_frame.paragraphs[0].font.bold = True
                    
                    workload_box = slide.shapes.add_textbox(Inches(5), Inches(y_pos), Inches(4), Inches(0.5))
                    workload_box.text_frame.text = f"{person.get('workload', 0)} tasks ({person.get('burden_level', 'Unknown')})"
                    workload_box.text_frame.paragraphs[0].font.size = Pt(14)
                    
                    y_pos += 0.7
    
    def _add_recommendations_slide(self, prs, data):
        """Recommendations slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "What To Do (1-2-3)"
        
        recs = data.get('recommendations', [])
        
        text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for rec in recs:
            p = tf.add_paragraph() if rec['priority'] > 1 else tf.paragraphs[0]
            p.text = f"{rec['priority']}. {rec['action']}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.space_after = Pt(6)
            
            why_p = tf.add_paragraph()
            why_p.text = f"   Why: {rec['why']}"
            why_p.font.size = Pt(14)
            why_p.space_after = Pt(12)
